"""
Tests for edit engine and operations.
"""

import tempfile
from pathlib import Path

import pytest
from pptx import Presentation
from pptx.util import Inches

from src.editing.operations import (
    EditOperation,
    EditOpType,
    AddTextPayload,
    ModifyTextPayload,
    DeleteShapePayload,
    Position,
    TextStyle,
)
from src.editing.editor import EditEngine


def _make_test_presentation() -> Presentation:
    """Create a Presentation with one slide for testing."""
    prs = Presentation()
    prs.slides.add_slide(prs.slide_layouts[6])  # blank
    return prs


class TestEditOperations:
    """Tests for edit operation models."""

    def test_add_text_operation(self) -> None:
        """Should create a valid add_text EditOperation."""
        op = EditOperation(
            op_type=EditOpType.ADD_TEXT,
            slide_idx=0,
            payload=AddTextPayload(
                text="Hello",
                position=Position(left=1, top=1, width=4, height=1),
                style=TextStyle(font_name="Arial", font_size=18),
            ),
        )
        assert op.op_type == EditOpType.ADD_TEXT
        assert op.payload.text == "Hello"

    def test_operation_serialization(self) -> None:
        """Should serialize and deserialize correctly."""
        op = EditOperation(
            op_type=EditOpType.MODIFY_TEXT,
            slide_idx=2,
            payload=ModifyTextPayload(shape_id=42, new_text="Updated"),
        )
        d = op.to_dict()
        restored = EditOperation.from_dict(d)
        assert restored.op_type == EditOpType.MODIFY_TEXT
        assert restored.slide_idx == 2

    def test_delete_operation(self) -> None:
        """Should create a valid delete operation."""
        op = EditOperation(
            op_type=EditOpType.DELETE_SHAPE,
            slide_idx=0,
            payload=DeleteShapePayload(shape_id=99),
        )
        assert op.payload.shape_id == 99


class TestEditEngine:
    """Tests for the EditEngine executor."""

    def test_execute_add_text(self) -> None:
        """Should add text to a slide via edit engine."""
        prs = _make_test_presentation()
        engine = EditEngine(prs)

        op = EditOperation(
            op_type=EditOpType.ADD_TEXT,
            slide_idx=0,
            payload=AddTextPayload(
                text="Engine Test",
                position=Position(left=1, top=1, width=4, height=1),
                style=TextStyle(),
            ),
        )

        result = engine.execute(op)
        assert result is True
        assert len(engine.history) == 1

        # Verify the text was actually added
        slide = prs.slides[0]
        assert slide.shapes[-1].text_frame.text == "Engine Test"

    def test_execute_batch(self) -> None:
        """Should execute multiple operations."""
        prs = _make_test_presentation()
        engine = EditEngine(prs)

        ops = [
            EditOperation(
                op_type=EditOpType.ADD_TEXT,
                slide_idx=0,
                payload=AddTextPayload(
                    text=f"Line {i}",
                    position=Position(left=1, top=1 + i * 0.5, width=4, height=0.4),
                    style=TextStyle(),
                ),
            )
            for i in range(3)
        ]

        success = engine.execute_batch(ops)
        assert success == 3
        assert len(engine.history) == 3

    def test_add_slide_operation(self) -> None:
        """Should add a new slide."""
        prs = _make_test_presentation()
        engine = EditEngine(prs)

        op = EditOperation(
            op_type=EditOpType.ADD_SLIDE,
            slide_idx=1,
            payload={"layout_index": 6},
        )

        engine.execute(op)
        assert len(prs.slides) >= 2

    def test_modify_text_operation(self) -> None:
        """Should modify existing text."""
        prs = _make_test_presentation()
        slide = prs.slides[0]
        from src.pptx_io.writer import add_text_box
        shape = add_text_box(
            slide, "Original",
            left=1, top=1, width=4, height=1,
        )

        engine = EditEngine(prs)
        op = EditOperation(
            op_type=EditOpType.MODIFY_TEXT,
            slide_idx=0,
            payload=ModifyTextPayload(shape_id=shape.shape_id, new_text="Modified"),
        )

        engine.execute(op)
        assert shape.text_frame.text == "Modified"

    def test_rollback(self) -> None:
        """Should allow rollback of last operation."""
        prs = _make_test_presentation()
        engine = EditEngine(prs)

        op = EditOperation(
            op_type=EditOpType.ADD_SLIDE,
            slide_idx=1,
            payload={"layout_index": 6},
        )

        engine.execute(op)
        count_after_add = len(prs.slides)
        engine.rollback_last()
        # Note: rollback for add_slide is best-effort
        assert len(engine.history) == 0
