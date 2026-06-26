"""
Tests for the evaluation module.
"""

import pytest
from pptx import Presentation

from src.evaluation.evaluator import (
    SlideEvaluator,
    RuleBasedEvaluator,
    EvaluationResult,
)


class TestRuleBasedEvaluator:
    """Tests for the rule-based (no-LLM) evaluator."""

    def test_evaluate_slide(self) -> None:
        """Should produce a score for a simple slide."""
        evaluator = RuleBasedEvaluator()

        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        from src.pptx_io.writer import add_text_box
        add_text_box(slide, "Test Title", left=1, top=0.5, width=8, height=1,
                     font_name="Arial", font_size=32, bold=True)
        add_text_box(slide, "Body text content here", left=1, top=2, width=8, height=3,
                     font_name="Arial", font_size=16)

        result = evaluator.evaluate(slide, 0)
        assert isinstance(result, EvaluationResult)
        assert 0 <= result.overall_score <= 10
        assert 0 <= result.content_richness <= 10
        assert 0 <= result.design_aesthetics <= 10
        assert 0 <= result.structural_coherence <= 10

    def test_empty_slide_scores_low(self) -> None:
        """Empty slides should get low scores."""
        evaluator = RuleBasedEvaluator()

        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

        result = evaluator.evaluate(slide, 0)
        # Empty slide should score poorly
        assert result.content_richness < 7.0
        assert result.overall_score < 7.0


class TestEvaluationResult:
    """Tests for EvaluationResult dataclass."""

    def test_serialization(self) -> None:
        """Should serialize to dict correctly."""
        result = EvaluationResult(
            slide_idx=0,
            content_richness=8.0,
            design_aesthetics=7.5,
            structural_coherence=9.0,
            overall_score=8.1,
            feedback="Looks good.",
            issues=["Minor alignment issue"],
            suggestions=["Align left margins"],
            is_acceptable=True,
        )
        d = result.to_dict()
        assert d["slide_idx"] == 0
        assert d["overall_score"] == 8.1
        assert len(d["issues"]) == 1
        assert d["is_acceptable"] is True

    def test_below_threshold(self) -> None:
        """Should correctly mark unacceptable scores."""
        result = EvaluationResult(
            slide_idx=0,
            content_richness=5.0,
            design_aesthetics=6.0,
            structural_coherence=5.5,
            overall_score=5.5,
            feedback="Needs work.",
            issues=["Too sparse"],
            suggestions=["Add more content"],
            is_acceptable=False,
        )
        assert not result.is_acceptable

        evaluator = SlideEvaluator(threshold=7.0)
        assert evaluator.needs_revision(result)
