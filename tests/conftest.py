"""
Shared pytest fixtures for PPT Agent tests.
"""

import tempfile
from pathlib import Path

import pytest
from pptx import Presentation
from pptx.util import Inches

from src.pptx_io.writer import add_text_box


@pytest.fixture
def sample_pptx() -> Path:
    """Create a temporary PPTX with 3 slides for testing."""
    prs = Presentation()

    # Slide 1 — title
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    add_text_box(slide1, "Title Slide", left=1, top=0.5, width=8, height=1.2,
                 font_name="Arial", font_size=36, bold=True)
    add_text_box(slide1, "Subtitle here", left=1, top=2, width=8, height=1,
                 font_name="Arial", font_size=20)

    # Slide 2 — content
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    add_text_box(slide2, "Content Title", left=0.5, top=0.3, width=9, height=0.8,
                 font_name="Arial", font_size=28, bold=True)
    add_text_box(slide2, "• Point 1\n• Point 2\n• Point 3", left=0.5, top=1.5, width=9, height=4,
                 font_name="Arial", font_size=16)

    # Slide 3 — section divider
    slide3 = prs.slides.add_slide(prs.slide_layouts[2])
    add_text_box(slide3, "SECTION BREAK", left=2, top=3, width=6, height=1.5,
                 font_name="Arial", font_size=40, bold=True)

    tmp = tempfile.NamedTemporaryFile(suffix=".pptx", delete=False)
    tmp_path = Path(tmp.name)
    prs.save(str(tmp_path))
    yield tmp_path
    tmp_path.unlink(missing_ok=True)


@pytest.fixture
def sample_text_file() -> Path:
    """Create a temporary text file with structured content."""
    content = """# Annual Report 2025

## Executive Summary

This is the executive summary for the annual report.
It covers the key highlights of the year.

## Financial Performance

Revenue increased by 15% compared to the previous year.
Operating margin improved to 22%.

| Quarter | Revenue | Profit |
|---------|---------|--------|
| Q1      | 100     | 20     |
| Q2      | 120     | 25     |
| Q3      | 140     | 30     |
| Q4      | 160     | 35     |

## Future Outlook

We expect continued growth in the coming year.
"""

    tmp = tempfile.NamedTemporaryFile(
        mode="w", suffix=".md", delete=False, encoding="utf-8"
    )
    tmp.write(content)
    tmp_path = Path(tmp.name)
    yield tmp_path
    tmp_path.unlink(missing_ok=True)


@pytest.fixture
def mock_llm_client():
    """Provide a mock LLM client for testing (to avoid real API calls)."""
    from unittest.mock import MagicMock
    mock = MagicMock()
    mock.model = "deepseek-chat"
    mock.chat.return_value = "Mock response"
    mock.json_chat.return_value = {"result": "ok"}
    return mock
