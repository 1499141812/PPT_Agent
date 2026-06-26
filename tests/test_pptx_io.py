"""
Tests for PPTX I/O — reader, writer, HTML converter.
"""

import tempfile
from pathlib import Path

import pytest
from pptx import Presentation
from pptx.util import Inches

from src.pptx_io.reader import (
    read_pptx,
    get_slide_count,
    extract_slide_shapes,
    duplicate_presentation,
)
from src.pptx_io.writer import (
    create_blank_presentation,
    add_text_box,
    add_slide,
    write_pptx,
)
from src.pptx_io.html_converter import (
    slide_to_html,
    pptx_to_html_snippets,
)


def _create_test_pptx() -> Path:
    """Create a minimal PPTX file for testing."""
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    add_text_box(
        slide,
        "Test Title",
        left=1.0, top=0.5, width=8.0, height=1.2,
        font_name="Arial",
        font_size=32,
        bold=True,
    )
    add_text_box(
        slide,
        "Subtitle text here",
        left=1.0, top=2.0, width=8.0, height=1.0,
        font_name="Arial",
        font_size=18,
    )

    tmp = tempfile.NamedTemporaryFile(suffix=".pptx", delete=False)
    tmp_path = Path(tmp.name)
    prs.save(str(tmp_path))
    return tmp_path


class TestPPTXReader:
    """Tests for PPT reading and shape extraction."""

    def test_read_pptx(self) -> None:
        """Should load a PPTX file."""
        tmp = _create_test_pptx()
        try:
            prs = read_pptx(tmp)
            assert prs is not None
            assert get_slide_count(prs) == 1
        finally:
            tmp.unlink()

    def test_extract_shapes(self) -> None:
        """Should extract shape metadata."""
        tmp = _create_test_pptx()
        try:
            prs = read_pptx(tmp)
            slide = prs.slides[0]
            shapes = extract_slide_shapes(slide)
            assert len(shapes) >= 2
            # One should be a text box
            text_shapes = [s for s in shapes if s.get("has_text")]
            assert len(text_shapes) >= 2
        finally:
            tmp.unlink()

    def test_duplicate_presentation(self) -> None:
        """Duplicated presentation should have same slide count."""
        tmp = _create_test_pptx()
        try:
            prs = read_pptx(tmp)
            dup = duplicate_presentation(prs)
            assert get_slide_count(dup) == get_slide_count(prs)
        finally:
            tmp.unlink()


class TestPPTXWriter:
    """Tests for PPT creation and writing."""

    def test_create_blank(self) -> None:
        """Should create a blank presentation."""
        prs = create_blank_presentation()
        assert prs is not None

    def test_add_slide_and_text(self) -> None:
        """Should add a slide with text."""
        prs = create_blank_presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout
        shape = add_text_box(
            slide,
            "Hello World",
            left=1.0, top=1.0, width=4.0, height=1.0,
        )
        assert shape is not None
        assert shape.has_text_frame
        assert shape.text_frame.text == "Hello World"

    def test_write_pptx(self) -> None:
        """Should save a PPTX to disk."""
        prs = create_blank_presentation()
        prs.slides.add_slide(prs.slide_layouts[6])
        with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as f:
            tmp_path = Path(f.name)
        try:
            write_pptx(prs, tmp_path)
            assert tmp_path.exists()
            assert tmp_path.stat().st_size > 0
        finally:
            tmp_path.unlink()


class TestHTMLConverter:
    """Tests for PPTX → HTML conversion."""

    def test_slide_to_html(self) -> None:
        """Should convert a slide to HTML."""
        tmp = _create_test_pptx()
        try:
            prs = read_pptx(tmp)
            html = slide_to_html(prs.slides[0], 0)
            assert '<div class="slide"' in html
            assert "Test Title" in html
            assert 'data-shape-type=' in html
        finally:
            tmp.unlink()

    def test_pptx_to_html_snippets(self) -> None:
        """Should convert all slides to HTML list."""
        tmp = _create_test_pptx()
        try:
            prs = read_pptx(tmp)
            snippets = pptx_to_html_snippets(prs)
            assert len(snippets) == 1
            assert "Test Title" in snippets[0]
        finally:
            tmp.unlink()
