"""
Main LangGraph workflow for PPT generation.

Node sequence:
    1. parse_documents    — parse source + reference PPT
    2. analyze_style      — ViT clustering + schema extraction
    3. plan_outline       — LLM generates slide-by-slide plan
    4. edit_loop          — per-slide: observe → edit → execute → evaluate → (revise?)
    5. finalize           — save PPT, build report
"""

from __future__ import annotations

import logging
import os
import shutil
from concurrent.futures import ThreadPoolExecutor, TimeoutError as FutureTimeoutError
from pathlib import Path
from typing import Any, Literal, Optional

from langgraph.graph import StateGraph, END

from src.config import get_config
from src.models import (
    AgentState,
    SlideState,
    create_initial_state,
)
from src.parsing import parse_document
from src.pptx_io import (
    read_pptx,
    get_slide_count,
    pptx_to_html_snippets,
    slide_to_html,
    create_blank_presentation,
    write_pptx,
)
from src.pptx_io.reader import duplicate_presentation
from src.pptx_io.slide_renderer import render_all_slides

logger = logging.getLogger(__name__)


# ── Workflow builder ────────────────────────────────────────────────────────

def build_workflow() -> StateGraph:
    """Build and return the LangGraph StateGraph for PPT generation.

    Returns:
        A compiled StateGraph ready to be invoked with an initial AgentState.
    """
    workflow = StateGraph(AgentState)

    # ── Add nodes ───────────────────────────────────────────────────────
    workflow.add_node("parse_documents", parse_documents_node)
    workflow.add_node("analyze_style", analyze_style_node)
    workflow.add_node("plan_outline", plan_outline_node)
    workflow.add_node("edit_slide", edit_slide_node)
    workflow.add_node("evaluate_slide", evaluate_slide_node)
    workflow.add_node("finalize", finalize_node)

    # ── Define edges ────────────────────────────────────────────────────
    workflow.set_entry_point("parse_documents")
    workflow.add_edge("parse_documents", "analyze_style")
    workflow.add_edge("analyze_style", "plan_outline")
    workflow.add_edge("plan_outline", "edit_slide")

    # After edit → evaluate
    workflow.add_edge("edit_slide", "evaluate_slide")

    # After evaluation: revise, next, or finish
    workflow.add_conditional_edges(
        "evaluate_slide",
        lambda s: (
            "finish" if s["current_slide_idx"] >= len(s["slide_states"])
            else s.get("_route", "next")
        ),
        {"revise": "edit_slide", "next": "edit_slide", "finish": "finalize"},
    )
    workflow.add_edge("finalize", END)

    return workflow.compile()


# ── Node implementations ────────────────────────────────────────────────────

def parse_documents_node(state: AgentState) -> AgentState:
    """Node 1: Parse source document and reference PPT.

    Inputs:  ``source_path``, ``reference_pptx_path``
    Outputs: ``source_doc``, ``reference_pptx``, ``reference_slide_count``
    """
    logger.info("=== Node 1: Parse Documents ===")

    # Parse source document
    source_path = state["source_path"]
    logger.info("Parsing source document: %s", source_path)
    source_doc = parse_document(source_path)
    state["source_doc"] = source_doc
    logger.info("Source parsed: %d sections, %d tables, %d images",
                len(source_doc["sections"]),
                len(source_doc["tables"]),
                len(source_doc["images"]))
    print(f"  [Step 1/6] 源文档: {len(source_doc['sections'])}章节, "
          f"{len(source_doc['tables'])}表格, {len(source_doc['images'])}图片")
    for img in source_doc.get("images", [])[:5]:
        p = img.get("path_to_saved_image", "?")
        print(f"    图片: {img.get('caption', '?')} → {p} "
              f"({'存在' if Path(p).exists() else '缺失'})")

    # Parse reference PPT
    ref_path = state["reference_pptx_path"]
    logger.info("Parsing reference PPT: %s", ref_path)
    ref_pptx = read_pptx(ref_path)
    state["reference_pptx"] = ref_pptx
    state["reference_slide_count"] = get_slide_count(ref_pptx)
    logger.info("Reference PPT has %d slides", state["reference_slide_count"])

    return state


def analyze_style_node(state: AgentState) -> AgentState:
    """Node 2: Analyze reference PPT style.

    1. Render slides to images
    2. Extract ViT features
    3. Hierarchical clustering
    4. Extract layout schemas via LLM

    **CRITICAL**: This entire node is wrapped in a hard timeout. If any sub-step
    hangs (e.g. ViT model download, LibreOffice subprocess), the timeout fires
    and style analysis is skipped. The workflow continues with empty schemas.
    """
    print("\n[Step 2/6] 分析参考PPT风格...")
    logger.info("=== Node 2: Analyze Style ===")

    cfg = get_config()

    ref_pptx = state["reference_pptx"]
    if ref_pptx is None or not cfg.style.style_analysis_enabled:
        reason = "无参考PPT" if ref_pptx is None else "STYLE_ANALYSIS_ENABLED=false"
        print(f"  ({reason}，跳过风格分析)")
        return state

    # ── Pre-flight check ─────────────────────────────────────────────────
    ready, messages = _check_style_readiness(cfg)
    for msg in messages:
        print(f"  {msg}")
    if not ready:
        print("  → 跳过风格分析，使用默认版式继续生成")
        return state

    # ── Run with hard timeout ─────────────────────────────────────────────
    # Runs _run_style_pipeline in a background thread. If it exceeds
    # STYLE_ANALYSIS_TIMEOUT seconds, the thread is abandoned and we
    # continue without style analysis.
    STYLE_TIMEOUT = int(os.getenv("STYLE_ANALYSIS_TIMEOUT", "300"))

    with ThreadPoolExecutor(max_workers=1) as executor:
        future = executor.submit(_run_style_pipeline, state, ref_pptx, cfg)
        try:
            future.result(timeout=STYLE_TIMEOUT)
        except FutureTimeoutError:
            print(f"  ⚠ 风格分析超过 {STYLE_TIMEOUT} 秒限制，自动跳过")
            print(f"  → 原因: ViT模型下载/HuggingFace连接可能卡住")
            print(f"  → 解决: 在 config.env 中设置 STYLE_ANALYSIS_ENABLED=false")
            print(f"  → 或设置 HF_ENDPOINT=https://hf-mirror.com 使用镜像")
            print(f"  → 生成将继续，使用默认版式")
            # Don't cancel the future — let the daemon thread die on exit
        except Exception as e:
            print(f"  ⚠ 风格分析失败: {e}")
            print(f"  → 生成将继续，使用默认版式")

    print()
    return state


def _run_style_pipeline(
    state: AgentState,
    ref_pptx: Any,
    cfg: Any,
) -> None:
    """Execute the full style analysis pipeline (called inside timeout)."""
    n_slides = len(ref_pptx.slides)

    # ── 2a. Render slides to images ─────────────────────────────────────
    print(f"  [2a] 渲染 {n_slides} 页幻灯片为图片...")
    try:
        if shutil.which("soffice"):
            print("     → 使用 LibreOffice 渲染")
        else:
            print("     → 使用 PIL 渲染 (纯 Python)")

        image_paths = render_all_slides(
            ref_pptx,
            cfg.temp_dir / "slide_images",
            dpi=cfg.style.slide_image_dpi,
        )
        print(f"     ✓ 渲染完成: {len(image_paths)} 张图片")
    except Exception as e:
        print(f"     ✗ 渲染失败: {e}")
        return

    if not image_paths:
        print("     ✗ 未生成任何图片，跳过后续分析")
        return

    # ── 2b. Extract ViT features ────────────────────────────────────────
    print(f"  [2b] 提取 ViT 视觉特征 (模型: {cfg.style.vit_model_name})...")
    from src.style.vit_extractor import ViTFeatureExtractor
    try:
        extractor = ViTFeatureExtractor()
        features = extractor.extract_features(image_paths)
        print(f"     ✓ 特征提取完成: {features.shape[0]} 页 x {features.shape[1]} 维")
    except Exception as e:
        print(f"     ✗ ViT 特征提取失败: {e}")
        return

    # ── 2c. Cluster slides ──────────────────────────────────────────────
    print(f"  [2c] 层次聚类 (识别版式分组)...")
    from src.style.clustering import cluster_slides
    clusters = cluster_slides(features)
    state["slide_clusters"] = clusters
    print(f"     ✓ 识别出 {len(clusters)} 种版式")
    for c in clusters:
        print(f"       - 版式 {c['cluster_id']}: {len(c['slide_indices'])} 页 "
              f"(代表页: #{c['representative_idx']})")

    # ── 2d. Extract layout schemas via LLM ──────────────────────────────
    print(f"  [2d] LLM 抽取版式 Schema...")
    html_snippets = pptx_to_html_snippets(ref_pptx)
    slide_html_map = {i: html for i, html in enumerate(html_snippets)}

    from src.style.schema_extractor import LayoutSchemaExtractor
    schema_extractor = LayoutSchemaExtractor()
    try:
        schemas = schema_extractor.extract_schemas(clusters, slide_html_map)

        # ── Extract style enrichment into a PLAIN DICT ──────────────────
        # CRITICAL: Store bg/fonts/shapes in state["style_enrichment"],
        # NOT inline in the TypedDict LayoutSchema. LangGraph may strip
        # extra keys from TypedDict fields during serialization between
        # nodes. A plain dict[str, Any] bypasses this.
        enrichment: dict[str, dict[str, Any]] = {}
        for s in schemas:
            sid = s["schema_id"]
            rep_idx = next(
                (c["representative_idx"] for c in clusters
                 if c["cluster_id"] == s["cluster_id"]),
                0,
            )
            if rep_idx < len(ref_pptx.slides):
                ref_slide = ref_pptx.slides[rep_idx]
                enrichment[sid] = {
                    "layout_index": _get_slide_layout_index(ref_pptx, ref_slide),
                    "slide_background_xml": _extract_slide_background(ref_slide),
                    "extracted_fonts": _extract_slide_fonts(ref_slide),
                    "shapes_xml": _extract_slide_shapes(ref_slide),
                }
            else:
                enrichment[sid] = {"layout_index": 0}

        state["layout_schemas"] = schemas
        state["style_enrichment"] = enrichment
        print(f"     ✓ 抽取完成: {len(schemas)} 个 Schema")
        for sid, ed in enrichment.items():
            n_shapes = len(ed.get("shapes_xml", []))
            has_bg = bool(ed.get("slide_background_xml"))
            fonts = ed.get("extracted_fonts", {})
            main_font = fonts.get("main_font", "?")
            print(f"       - {sid}: layout {ed.get('layout_index', '?')}, "
                  f"BG={'✓' if has_bg else '✗'}, "
                  f"shapes={n_shapes}, "
                  f"font={main_font})")
    except Exception as e:
        print(f"     ✗ Schema 抽取失败: {e}")
        print(f"     → 将使用空 Schema 继续，由 LLM 自行决定布局")


def _get_background_for_slide(
    slide_state: dict[str, Any],
    schemas: list[dict[str, Any]],
) -> str | None:
    """Get the per-slide background XML for a slide's schema."""
    for s in schemas:
        if s["schema_id"] == slide_state["schema_id"]:
            return s.get("slide_background_xml")
    return None


def _prepare_slide_for_editing(slide: Any) -> dict[int, dict]:
    """Clear text/remove content shapes. Returns saved font properties."""
    from src.pptx_io.writer import clear_slide_text

    for shape in list(slide.shapes):
        try:
            stype = str(shape.shape_type) if shape.shape_type else ''
        except Exception:
            stype = ''
        if 'TABLE' in stype or 'CHART' in stype:
            sp = shape._element
            sp.getparent().remove(sp)

    return clear_slide_text(slide)


def _get_rep_slide_index(
    slide_state: dict[str, Any],
    clusters: list[dict[str, Any]],
) -> int:
    """Find the representative slide index for the current slide's cluster."""
    for c in clusters:
        if c["cluster_id"] == slide_state["cluster_id"]:
            return c.get("representative_idx", 0)
    return 0


def _closing_text_from_outline(outline: list) -> str:
    """Generate closing text based on the presentation topic."""
    title = outline[0].get("title", "") if outline else ""
    if "VS" in title.upper() or "对比" in title or "差异" in title:
        return "感谢聆听"
    return "Thank You  &  Q&A"


def _generate_extra_outline_items(
    source_doc: dict[str, Any],
    start_idx: int,
    target_count: int,
) -> list[dict[str, Any]]:
    """Generate minimal outline items to pad a short plan.

    Creates simple title+bullet slides from document sections that
    weren't covered by the LLM's plan.
    """
    items: list[dict[str, Any]] = []
    sections = source_doc.get("sections", [])
    for i in range(start_idx, target_count):
        sec = sections[i % len(sections)] if sections else {}
        items.append({
            "slide_idx": i,
            "title": sec.get("heading", f"Slide {i + 1}"),
            "narrative_role": "content",
            "cluster_id": 0,
            "schema_id": "schema_0",
            "content": {
                "bullet_points": [
                    p[:200] for p in sec.get("paragraphs", [])[:5]
                ] or [f"Content for slide {i + 1}"],
            },
            "content_summary": sec.get("heading", ""),
        })
    return items


def _get_layout_for_slide(
    slide_state: dict[str, Any],
    schemas: list[dict[str, Any]],
) -> int:
    """Get the appropriate slide layout index for a new slide.

    Uses the layout recorded from the reference PPT's representative
    slide for this schema/cluster. Falls back to layout 0 (usually
    the main layout with theme background).
    """
    for s in schemas:
        if s["schema_id"] == slide_state["schema_id"]:
            layout = s.get("layout_index")
            if layout is not None:
                return layout
    return 0  # Default: first layout (has theme background)


def _extract_slide_fonts(slide: Any) -> dict[str, Any]:
    """Extract the dominant font styles from a reference slide."""
    from collections import Counter

    fonts: list[str] = []
    sizes: list[float] = []
    colors: list[str] = []

    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for para in shape.text_frame.paragraphs:
            # Check paragraph-level font (defRPr — often the primary source)
            pf = para.font
            if pf.name:
                fonts.append(pf.name)
            if pf.size:
                sizes.append(pf.size / 12700)
            try:
                if pf.color and pf.color.rgb:
                    colors.append(str(pf.color.rgb))
            except Exception:
                pass

            # Check run-level fonts (overrides)
            for run in para.runs:
                rf = run.font
                name = rf.name or pf.name  # fall back to paragraph font
                if name:
                    fonts.append(name)
                size = rf.size or pf.size
                if size:
                    sizes.append(size / 12700)
                try:
                    if rf.color and rf.color.rgb:
                        colors.append(str(rf.color.rgb))
                    elif pf.color and pf.color.rgb:
                        colors.append(str(pf.color.rgb))
                except Exception:
                    pass

    font_counts = Counter(fonts)
    size_counts = Counter(sizes)
    color_counts = Counter(colors)

    main_font = font_counts.most_common(1)[0][0] if font_counts else "Arial"
    # Cap sizes — 90pt from reference is for its specific design,
    # using it on new slides would overflow
    raw_title = max(sizes) if sizes else 32
    raw_body = size_counts.most_common(1)[0][0] if size_counts else 16
    title_size = min(raw_title, 44)
    body_size = min(raw_body, 20)
    main_color = color_counts.most_common(1)[0][0] if color_counts else "333333"

    return {
        "main_font": main_font,
        "title_size": int(title_size),
        "body_size": int(body_size),
        "main_color": main_color,
        "all_fonts": [f for f, _ in font_counts.most_common(5)],
        "all_colors": [c for c, _ in color_counts.most_common(5)],
    }


def _extract_slide_shapes(slide: Any) -> list[dict[str, Any]]:
    """Extract all shapes from a reference slide (pure template — keep everything)."""
    from pptx.oxml.ns import qn
    from lxml import etree

    # Slide dimensions from the presentation, not the slide itself
    prs = slide.part.package.presentation_part.presentation
    slide_w = prs.slide_width or 9144000
    slide_h = prs.slide_height or 5143500
    slide_area = slide_w * slide_h

    shapes: list[dict[str, Any]] = []
    temp_dir = get_config().temp_dir

    for shape in slide.shapes:
        stype = str(shape.shape_type) if shape.shape_type else 'unknown'

        # ── Picture: extract image blob for re-adding ────────────────
        if 'PICTURE' in stype or 'PICTURE' in str(type(shape).__name__).upper():
            try:
                blob = shape.image.blob
                ext = shape.image.content_type.split('/')[-1]
                img_path = temp_dir / f"ref_img_{shape.shape_id}.{ext}"
                img_path.write_bytes(blob)
                shapes.append({
                    "type": "picture",
                    "image_path": str(img_path),
                    "left": shape.left,
                    "top": shape.top,
                    "width": shape.width,
                    "height": shape.height,
                })
            except Exception:
                pass  # Can't extract image — skip
            continue

        # ── Non-picture: use XML extraction ──────────────────────────
        child = shape._element
        tag = child.tag.split('}')[-1] if '}' in child.tag else child.tag
        if tag not in ('sp', 'graphicFrame', 'grpSp', 'cxnSp'):
            continue

        xml_str = etree.tostring(child, encoding='unicode')
        shapes.append({"type": "xml", "xml": xml_str})

    return shapes


def _inject_shapes(slide: Any, shapes: list[dict[str, Any]]) -> None:
    """Inject decorative shapes into a slide.

    Handles two types:
      - {"type": "xml", "xml": "<...>"}     → inject XML (auto-shapes)
      - {"type": "picture", "image_path": …, "left": …, ...} → add via add_picture
    """
    if not shapes:
        return
    from lxml import etree
    from pptx.oxml.ns import qn
    from pptx.util import Inches

    spTree = slide.shapes._spTree
    for item in shapes:
        try:
            # ── Picture type: add via python-pptx (handles relationships) ──
            if item.get("type") == "picture":
                img_path = item.get("image_path", "")
                if img_path:
                    slide.shapes.add_picture(
                        img_path,
                        item["left"], item["top"],
                        item.get("width"), item.get("height"),
                    )
                continue

            # ── XML type: inject shape XML ─────────────────────────────
            xml_str = item.get("xml", "")
            if not xml_str:
                continue

            elem = etree.fromstring(xml_str.encode('utf-8'))
            tag = elem.tag.split('}')[-1] if '}' in elem.tag else elem.tag

            # Handle fills
            spPr = elem.find(qn('p:spPr'))
            if spPr is not None:
                blip_fill = spPr.find(qn('a:blipFill'))
                solid_fill = spPr.find(qn('a:solidFill'))
                grad_fill = spPr.find(qn('a:gradFill'))

                # If shape has blipFill (picture fill) AND a solid/gradient fallback,
                # remove the blipFill so the shape renders with its solid color.
                # If ONLY blipFill, skip the shape entirely.
                if blip_fill is not None:
                    if solid_fill is not None or grad_fill is not None:
                        spPr.remove(blip_fill)
                    else:
                        continue

                # Strip theme-dependent colors (schemeClr → srgbClr)
                for fill in (solid_fill, grad_fill):
                    if fill is None:
                        continue
                    for scheme_clr in fill.findall('.//' + qn('a:schemeClr')):
                        parent_s = scheme_clr.getparent()
                        if parent_s is not None:
                            parent_s.remove(scheme_clr)
                            new_srgb = etree.SubElement(parent_s, qn('a:srgbClr'))
                            new_srgb.set('val', 'FFFFFF')

            spTree.append(elem)
        except Exception:
            pass


def _enforce_font_style(slide: Any, fonts: dict[str, Any]) -> None:
    """Apply font name/color + restore per-shape sizes after clearing."""
    from pptx.util import Pt
    from src.pptx_io.writer import _hex_to_rgb

    main_font = fonts.get("main_font", "")
    main_color = fonts.get("main_color", "")
    saved = fonts.get("_saved_fonts", {})

    if not main_font and not main_color and not saved:
        return

    rgb = _hex_to_rgb(main_color) if main_color else None

    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        sid = shape.shape_id
        sf = saved.get(sid, {})

        for para in shape.text_frame.paragraphs:
            try:
                if main_font:
                    para.font.name = main_font
            except Exception: pass
            try:
                if rgb:
                    para.font.color.rgb = rgb
            except Exception: pass
            try:
                sz = sf.get("font_size")
                if sz:
                    para.font.size = Pt(sz)
            except Exception: pass
            try:
                if sf.get("bold") is not None:
                    para.font.bold = sf["bold"]
            except Exception: pass

            for run in para.runs:
                try:
                    if main_font: run.font.name = main_font
                except Exception: pass
                try:
                    if rgb: run.font.color.rgb = rgb
                except Exception: pass
                try:
                    if sz: run.font.size = Pt(sz)
                except Exception: pass
                try:
                    if sf.get("bold") is not None: run.font.bold = sf["bold"]
                except Exception: pass


def _get_slide_layout_index(prs: Any, slide: Any) -> int:
    """Find which layout index a slide uses in its presentation."""
    slide_layout = slide.slide_layout
    for idx, layout in enumerate(prs.slide_layouts):
        if layout is slide_layout:
            return idx
    return 0


def _extract_slide_background(slide: Any) -> str | None:
    """Extract per-slide background XML from a reference slide.

    Returns the serialized <p:bg> element if present, or None.
    This is needed because many templates set background PER SLIDE,
    not on the master/layout.
    """
    from pptx.oxml.ns import qn
    from lxml import etree

    cSld = slide._element.find(qn('p:cSld'))
    if cSld is None:
        return None
    bg = cSld.find(qn('p:bg'))
    if bg is None:
        return None
    return etree.tostring(bg, encoding='unicode')


def _apply_slide_background(slide: Any, bg_xml: str | None) -> None:
    """Apply a per-slide background to a newly created slide.

    Extracts the fill color/type from the reference XML and creates a
    fresh <p:bg> element on the new slide. This avoids relationship
    issues that would occur from copying XML directly.
    """
    if not bg_xml:
        return

    import copy
    from pptx.oxml.ns import qn
    from lxml import etree

    # Parse the reference background
    try:
        bg_elem = etree.fromstring(bg_xml.encode('utf-8'))
    except Exception:
        return

    bgPr = bg_elem.find(qn('p:bgPr'))
    if bgPr is None:
        return

    # Check fill type
    solid = bgPr.find(qn('a:solidFill'))
    grad = bgPr.find(qn('a:gradFill'))
    noFill = bgPr.find(qn('a:noFill'))

    # Build new bgPr element
    new_bgPr = etree.Element(qn('p:bgPr'))

    if noFill is not None:
        etree.SubElement(new_bgPr, qn('a:noFill'))
    elif solid is not None:
        # Copy the solid fill (uses inline colors, no external refs)
        new_bgPr.append(copy.deepcopy(solid))
    elif grad is not None:
        # Copy the gradient (uses inline colors)
        new_bgPr.append(copy.deepcopy(grad))
    else:
        return  # Unknown fill type — skip

    # Create new <p:bg> with the bgPr
    new_bg = etree.Element(qn('p:bg'))
    new_bg.append(new_bgPr)

    # Get or create cSld
    cSld = slide._element.find(qn('p:cSld'))
    if cSld is None:
        cSld = etree.SubElement(slide._element, qn('p:cSld'))

    # Remove any existing background
    existing = cSld.find(qn('p:bg'))
    if existing is not None:
        cSld.remove(existing)

    # Insert at the beginning
    cSld.insert(0, new_bg)


# ── Pre-flight check for style analysis ─────────────────────────────────────

def _check_style_readiness(cfg: Any) -> tuple[bool, list[str]]:
    """Quick checks before attempting style analysis.

    Returns:
        (is_ready, messages) — whether to proceed and diagnostic messages.
    """
    messages: list[str] = []
    ready = True

    # Check 1: Is the ViT model already cached locally?
    vit_model = cfg.style.vit_model_name
    from huggingface_hub import try_to_load_from_cache
    try:
        cached = try_to_load_from_cache(
            repo_id=vit_model,
            filename="pytorch_model.bin",
        ) or try_to_load_from_cache(
            repo_id=vit_model,
            filename="model.safetensors",
        )
        if cached:
            messages.append(f"  ✓ ViT 模型已缓存")
        else:
            messages.append(f"  ⚠ ViT 模型未缓存，首次需下载 ~350MB")
            messages.append(f"    下载地址: https://huggingface.co/{vit_model}")
            messages.append(f"    使用国内镜像: 在 config.env 中设置 HF_ENDPOINT=https://hf-mirror.com")
            # Don't block — let the timeout handle it, but warn
    except ImportError:
        # huggingface_hub not available for cache check — proceed anyway
        pass
    except Exception:
        pass

    # Check 2: Is LibreOffice responsive? (quick 5s test)
    if shutil.which("soffice"):
        try:
            import subprocess
            result = subprocess.run(
                ["soffice", "--version"],
                capture_output=True, text=True, timeout=5,
            )
            if result.returncode == 0:
                messages.append(f"  ✓ LibreOffice 可用: {result.stdout.strip()[:60]}")
            else:
                messages.append(f"  ⚠ LibreOffice 异常，将使用 PIL 渲染")
        except subprocess.TimeoutExpired:
            messages.append(f"  ⚠ LibreOffice 无响应 (5s超时)，将使用 PIL 渲染")
        except Exception:
            messages.append(f"  ⚠ LibreOffice 检测失败，将使用 PIL 渲染")

    return ready, messages


def plan_outline_node(state: AgentState) -> AgentState:
    """Node 3: Generate slide-by-slide outline.

    Inputs:  ``source_doc``, ``layout_schemas``
    Outputs: ``outline``, initializes ``slide_states`` and creates output PPTX
    """
    logger.info("=== Node 3: Plan Outline ===")

    source_doc = state["source_doc"]
    schemas = state["layout_schemas"]

    if source_doc is None:
        logger.warning("No source document — cannot plan outline")
        return state

    from src.planning.outline_planner import OutlinePlanner
    planner = OutlinePlanner()

    # Determine min slides based on source content
    n_sections = len(source_doc.get("sections", []))
    min_slides = max(n_sections, 3)  # at least title + content + ending

    outline = planner.plan(source_doc, schemas or [])

    # Ensure minimum slide count
    if len(outline) < min_slides:
        print(f"  ⚠ LLM 只规划了 {len(outline)} 页 (源文档有 {n_sections} 章节)")
        print(f"     → 补足到最少 {min_slides} 页")
        # Pad with extra slides from remaining sections
        extra = _generate_extra_outline_items(source_doc, len(outline), min_slides)
        outline.extend(extra)

    state["outline"] = outline
    print(f"  [Step 3/6] 大纲规划: {len(outline)} 页")
    logger.info("Planned %d slides", len(outline))

    # ── Initialize the output presentation ──────────────────────────────
    # Keep the reference PPT INTACT with all its slides. We'll clone
    # representative slides in edit_slide_node. This preserves ALL shapes,
    # fills, images, and decorations without any XML manipulation.
    ref_pptx = state.get("reference_pptx")
    if ref_pptx is not None:
        from src.pptx_io.reader import duplicate_presentation
        state["output_pptx"] = duplicate_presentation(ref_pptx)
        # Map cluster_id → slide index of representative in the output PPT
        for c in state.get("slide_clusters", []):
            c["_rep_slide_idx"] = c.get("representative_idx", 0)
        logger.info("Output PPTX initialized: %d template slides",
                    len(state["output_pptx"].slides))
    else:
        state["output_pptx"] = create_blank_presentation()
        logger.info("Output PPTX initialized as blank presentation")

    # ── Initialize slide states ─────────────────────────────────────────
    # Each slide starts with a truly empty HTML view matching a blank slide.
    EMPTY_HTML = (
        '<div class="slide" id="slide-{idx}" '
        'style="width:10.00in;height:7.50in;position:relative;overflow:hidden;background:#ffffff;">'
        '</div>'
    )
    slide_states: list[SlideState] = []
    for item in outline:
        slide_states.append(SlideState(
            slide_idx=item["slide_idx"],
            cluster_id=item.get("cluster_id", 0),
            schema_id=item.get("schema_id", ""),
            html_view=EMPTY_HTML.format(idx=item["slide_idx"]),
            edit_history=[],
            evaluation_scores={},
            overall_score=0.0,
            revision_round=0,
            is_acceptable=False,
        ))
    state["slide_states"] = slide_states
    state["current_slide_idx"] = 0

    return state


def edit_slide_node(state: AgentState) -> AgentState:
    """Node 4: Generate and execute edits for the current slide.

    Uses SINGLE-TURN JSON mode (fast, one LLM call per slide).
    Every slide starts blank — always ADD-only, no modify/delete needed.
    """
    current_idx = state["current_slide_idx"]
    # Guard: stop if past last slide (evaluate advances past end)
    if current_idx >= len(state["slide_states"]):
        return state

    outline = state["outline"]
    slide_state = state["slide_states"][current_idx]

    print(f"\n  [Slide {current_idx + 1}/{len(outline)}] "
          f"生成: {outline[current_idx].get('title', 'Untitled')}")

    schemas = state["layout_schemas"]
    source_doc = state["source_doc"]
    output_pptx = state["output_pptx"]

    # Find schema + enrichment for this slide
    schema: dict[str, Any] = {}
    enrichment: dict[str, Any] = {}
    for s in schemas:
        if s["schema_id"] == slide_state["schema_id"]:
            schema = s
            break
    # Read enrichment from the PLAIN DICT (safe from TypedDict stripping)
    enrichment = state.get("style_enrichment", {}).get(slide_state["schema_id"], {})

    # Diagnostic: show what style data was extracted
    if current_idx == 0:
        has_bg = bool(enrichment.get("slide_background_xml"))
        n_shapes = len(enrichment.get("shapes_xml", []))
        fonts = enrichment.get("extracted_fonts", {})
        print(f"     [风格数据] BG={'✓' if has_bg else '✗'} shapes={n_shapes} "
              f"font={fonts.get('main_font', '?')} "
              f"title_size={fonts.get('title_size', '?')} "
              f"color=#{fonts.get('main_color', '?')}")

    # Ensure the slide exists and is ready for editing
    # For extra slides beyond template count, duplicate a random template
    import random
    is_new_slide = current_idx >= len(output_pptx.slides)
    while current_idx >= len(output_pptx.slides):
        from src.pptx_io.writer import duplicate_slide, clear_slide_text
        if len(output_pptx.slides) > 0:
            src = random.choice(list(output_pptx.slides))
            new_s = duplicate_slide(output_pptx, src)
            clear_slide_text(new_s)  # Clear old text, keep decorations + images
        else:
            from src.pptx_io.writer import add_slide
            add_slide(output_pptx, layout_index=0)

    actual_slide = output_pptx.slides[current_idx]

    # ALWAYS clear text on existing slides (template or new)
    from src.pptx_io.writer import clear_slide_text
    saved_fonts = _prepare_slide_for_editing(actual_slide)
    if saved_fonts:
        enrichment = dict(enrichment)
        enrichment["_saved_fonts"] = saved_fonts

    # Apply background to new blank slides
    if is_new_slide:
        bg_xml = enrichment.get("slide_background_xml")
        if bg_xml:
            _apply_slide_background(actual_slide, bg_xml)

    # Verify text cleared
    leftover = sum(1 for s in actual_slide.shapes
                   if s.has_text_frame and s.text_frame.text.strip())
    if leftover > 0 and current_idx < 3:
        print(f"     ⚠ {leftover} shapes still have text: "
              + ", ".join(f"'{s.text_frame.text[:30]}'" for s in actual_slide.shapes
                          if s.has_text_frame and s.text_frame.text.strip()))

    # Content for this slide
    content: dict[str, Any] = {}
    if current_idx < len(outline):
        content = outline[current_idx].get("content", {})
    # On revision, inject previous feedback
    if slide_state["revision_round"] > 0:
        fb = slide_state.get("evaluation_feedback", "")
        if fb:
            content["_revision_feedback"] = fb

    # Images: only real existing files
    raw_images = source_doc.get("images", []) if source_doc else []
    available_images = [
        img for img in raw_images
        if img.get("path_to_saved_image")
        and Path(img["path_to_saved_image"]).exists()
    ]
    available_tables = source_doc.get("tables", []) if source_doc else []
    if current_idx == 0 and available_images:
        print(f"    源文档图片: {len(available_images)}张可用 "
              f"({len(raw_images)}张提取, {len(raw_images) - len(available_images)}张缺失)")

    # Get HTML view of actual slide
    actual_slide = output_pptx.slides[current_idx]
    slide_state["html_view"] = slide_to_html(actual_slide, current_idx)

    # ── Single LLM call ───────────────────────────────────────────────
    print(f"     → LLM 生成编辑指令...")
    from src.editing.llm_editor import LLMEditor
    llm_editor = LLMEditor()

    # Inject font info + shape sizes into content for LLM
    fonts = enrichment.get("extracted_fonts", {})
    if fonts:
        content["_fonts"] = fonts
    if enrichment.get("_saved_fonts"):
        content["_shape_sizes"] = enrichment["_saved_fonts"]

    # Determine if slide has existing shapes (cloned from template)
    # vs. being truly blank (added as new). Cloned slides have shapes
    # that should be filled via modify_text, not add_text.
    has_existing_shapes = len(actual_slide.shapes) > 2

    try:
        edit_ops = llm_editor.generate_edits_single_turn(
            slide_html=slide_state["html_view"],
            schema=schema,
            content=content,
            is_fresh_slide=not has_existing_shapes,
            available_images=available_images,
            available_tables=available_tables,
        )
    except Exception as e:
        logger.warning("LLM edit generation failed for slide %d: %s", current_idx, e)
        print(f"     ⚠ LLM 调用失败: {e}")
        edit_ops = []

    logger.info("Generated %d edits for slide %d", len(edit_ops), current_idx)
    if not edit_ops:
        print(f"     ⚠ LLM 返回空操作列表 — 跳过")

    # ── Diagnostic: show what was generated ────────────────────────────
    for op in edit_ops[:5]:
        stype = op.op_type.value if hasattr(op, 'op_type') else str(op)
        if hasattr(op, 'payload'):
            sid = getattr(op.payload, 'shape_id', 'N/A')
            txt = getattr(op.payload, 'new_text', getattr(op.payload, 'text', ''))[:40]
            print(f"     [{stype}] shape_id={sid} text='{txt}'")

    # ── Fix slide_idx ───────────────────────────────────────────────────
    for op in edit_ops:
        if hasattr(op, 'slide_idx'):
            op.slide_idx = current_idx

    # ── Diagnostic: show actual shape IDs on the slide ──────────────────
    actual_ids = [s.shape_id for s in actual_slide.shapes if s.has_text_frame]
    if current_idx == 0:
        print(f"     Slide shapes with text: {actual_ids}")

    # ── Execute with per-op diagnostics ─────────────────────────────────
    from src.editing.editor import EditEngine
    engine = EditEngine(output_pptx)
    for op in edit_ops:
        try:
            engine.execute(op)
        except Exception as e:
            print(f"     ✗ {op.op_type.value}: {str(e)[:80]}")
    success = sum(1 for op in edit_ops if True)  # approximate
    if len(edit_ops) > 0:
        final_text = sum(1 for s in actual_slide.shapes if s.has_text_frame and s.text_frame.text.strip())
        print(f"     → 最终 {final_text} 个形状有文字")

    # ── Post-edit: enforce fonts with per-shape saved sizes ────────────
    enforce_data = dict(fonts) if fonts else {}
    if enrichment.get("_saved_fonts"):
        enforce_data["_saved_fonts"] = enrichment["_saved_fonts"]
    if enforce_data:
        _enforce_font_style(actual_slide, enforce_data)

    # Update state
    slide_state["edit_history"].extend(
        op.to_dict() if hasattr(op, "to_dict") else op for op in edit_ops
    )
    if current_idx < len(output_pptx.slides):
        slide_state["html_view"] = slide_to_html(output_pptx.slides[current_idx], current_idx)
    state["edit_log"].extend(
        op.to_dict() if hasattr(op, "to_dict") else op for op in edit_ops
    )

    return state


def evaluate_slide_node(state: AgentState) -> AgentState:
    """Node 5: LLM evaluate + trigger revision if score < threshold."""
    current_idx = state["current_slide_idx"]
    if current_idx >= len(state["slide_states"]):
        return state

    slide_state = state["slide_states"][current_idx]
    output_pptx = state["output_pptx"]
    MAX_REVISIONS = 2

    if current_idx < len(output_pptx.slides):
        html = slide_to_html(output_pptx.slides[current_idx], current_idx)
        outline = state.get("outline", [])
        summary = outline[current_idx].get("content_summary", "") if current_idx < len(outline) else ""

        # LLM evaluation
        from src.evaluation.evaluator import llm_evaluate
        score, suggestions = llm_evaluate(html, summary)
        slide_state["overall_score"] = score
        slide_state["is_acceptable"] = score >= 6.0

        if suggestions:
            slide_state["evaluation_feedback"] = suggestions
            slide_state["evaluation_suggestions"] = suggestions.split("\n") if suggestions else []

        print(f"     评分: {score:.1f}/10", end="")
        if suggestions:
            print(f" — {suggestions[:80]}...")
        else:
            print()

        # Revision logic
        if score < 6.0 and slide_state["revision_round"] < MAX_REVISIONS:
            slide_state["revision_round"] += 1
            state["_route"] = "revise"
            print(f"     🔄 第{slide_state['revision_round']}/{MAX_REVISIONS}次修正")
        else:
            state["current_slide_idx"] = current_idx + 1
            state["_route"] = "next"
    else:
        state["current_slide_idx"] = current_idx + 1
        state["_route"] = "next"

    return state


def finalize_node(state: AgentState) -> AgentState:
    """Node 6: Save the final PPTX and generate a summary report.

    Inputs:  ``output_pptx``, ``output_pptx_path``
    Outputs: ``final_report``
    """
    logger.info("=== Node 6: Finalize ===")

    output_pptx = state["output_pptx"]
    output_path = state.get("output_pptx_path", "output/generated.pptx")
    outline_len = len(state.get("outline", []))

    if output_pptx is not None:
        from src.pptx_io.writer import delete_slide

        # Delete extra slides, but keep the LAST one as closing (Thank You / Q&A)
        while len(output_pptx.slides) > outline_len + 1:
            delete_slide(output_pptx, len(output_pptx.slides) - 2)  # second-to-last

        # If there's exactly one extra slide, keep it as closing
        if len(output_pptx.slides) == outline_len + 1:
            pass  # Keep the extra slide as closing

        # If closing slide was preserved, clear its text and add ending
        if len(output_pptx.slides) > outline_len:
            closing = output_pptx.slides[-1]
            from src.pptx_io.writer import clear_slide_text, add_text_box
            clear_slide_text(closing)
            closing_text = _closing_text_from_outline(state.get("outline", []))
            # Find largest text box or create one
            best_shape = None
            best_size = 0
            for s in closing.shapes:
                if s.has_text_frame:
                    try:
                        p0 = s.text_frame.paragraphs[0]
                        r0 = p0.runs[0] if p0.runs else None
                        sz = r0.font.size if r0 else p0.font.size
                        if sz and (sz > best_size):
                            best_size = sz
                            best_shape = s
                    except Exception:
                        pass
            if best_shape and best_shape.has_text_frame:
                from src.pptx_io.writer import modify_text
                modify_text(best_shape, closing_text)
            else:
                # No text shape found — add one at bottom, avoiding images
                from pptx.util import Inches
                add_text_box(closing, closing_text,
                             left=1.0, top=5.5, width=8.0, height=1.0,
                             font_size=28, bold=True)
            logger.info("Closing slide text set to: %s", closing_text)

        logger.info("Final slide count: %d (outline: %d, closing preserved)",
                    len(output_pptx.slides), outline_len)

        write_pptx(output_pptx, output_path)
        logger.info("PPT saved to: %s", output_path)
    else:
        logger.warning("No output PPTX to save!")

    # Build report
    slide_states = state.get("slide_states", [])
    total_slides = len(slide_states)
    avg_score = (
        sum(s["overall_score"] for s in slide_states) / total_slides
        if total_slides > 0 else 0.0
    )
    revision_count = sum(s["revision_round"] for s in slide_states)
    total_edits = len(state.get("edit_log", []))

    report = {
        "total_slides": total_slides,
        "average_score": round(avg_score, 2),
        "total_revisions": revision_count,
        "total_edit_operations": total_edits,
        "per_slide_scores": [
            {
                "slide_idx": s["slide_idx"],
                "score": s["overall_score"],
                "revisions": s["revision_round"],
                "acceptable": s["is_acceptable"],
            }
            for s in slide_states
        ],
        "output_path": str(output_path),
    }

    state["final_report"] = report

    logger.info("=== Generation Complete ===")
    logger.info("Slides: %d | Avg Score: %.2f | Revisions: %d | Edits: %d",
                total_slides, avg_score, revision_count, total_edits)

    return state


# ── Conditional routing ─────────────────────────────────────────────────────

# ── Convenience runner ──────────────────────────────────────────────────────

def run_ppt_generation(
    source_path: str | Path,
    reference_pptx_path: str | Path,
    output_pptx_path: str | Path = "output/generated.pptx",
) -> AgentState:
    """Run the complete PPT generation workflow.

    This is the primary entry point for programmatic use.

    Args:
        source_path: Path to the source document (.docx, .pdf, or .txt).
        reference_pptx_path: Path to the reference PPT for style extraction.
        output_pptx_path: Where to save the generated PPT.

    Returns:
        The final ``AgentState`` with all results.
    """
    app = build_workflow()

    initial_state = create_initial_state(
        source_path=str(source_path),
        reference_pptx_path=str(reference_pptx_path),
        output_pptx_path=str(output_pptx_path),
    )

    final_state = app.invoke(initial_state)
    return final_state
