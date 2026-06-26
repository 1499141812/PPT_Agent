"""
PPTX reader — loads reference presentations and extracts structural data.
"""

from __future__ import annotations

import copy
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.shapes import MSO_SHAPE_TYPE


def read_pptx(file_path: str | Path) -> Presentation:
    """Load a .pptx file and return the Presentation object.

    Args:
        file_path: Path to the reference PPT.

    Returns:
        A ``pptx.Presentation`` instance.

    Raises:
        FileNotFoundError: If the file does not exist.
    """
    path = Path(file_path).resolve()
    if not path.exists():
        raise FileNotFoundError(f"PPT file not found: {path}")
    return Presentation(str(path))


def get_slide_count(pptx: Presentation) -> int:
    """Return the number of slides in a presentation."""
    return len(pptx.slides)


def extract_slide_shapes(slide: Any) -> list[dict[str, Any]]:
    """Extract metadata for every shape on a slide.

    Each shape dict contains::

        {
            "shape_id": int,
            "name": str,
            "shape_type": str,      # "text_box", "picture", "table", "chart", "group", "auto_shape"
            "left": int (EMU), "top": int (EMU),
            "width": int (EMU), "height": int (EMU),
            "text": str | None,
            "font_name": str | None,
            "font_size": float | None,   # in points
            "bold": bool,
            "italic": bool,
            "color": str | None,         # hex
            "fill_color": str | None,
            "image_blob": bytes | None,  # only for pictures
            "has_text": bool,
            "placeholder_type": str | None,
        }

    Args:
        slide: A ``pptx.slide.Slide`` object.

    Returns:
        List of shape metadata dicts.
    """
    shapes: list[dict[str, Any]] = []
    for shape in slide.shapes:
        info = _extract_shape_info(shape)
        shapes.append(info)
    return shapes


def _extract_shape_info(shape: Any) -> dict[str, Any]:
    """Extract metadata from a single shape."""
    info: dict[str, Any] = {
        "shape_id": shape.shape_id,
        "name": shape.name,
        "shape_type": str(shape.shape_type) if shape.shape_type else "unknown",
        "left": shape.left,
        "top": shape.top,
        "width": shape.width,
        "height": shape.height,
        "text": None,
        "font_name": None,
        "font_size": None,
        "bold": False,
        "italic": False,
        "color": None,
        "fill_color": None,
        "image_blob": None,
        "has_text": shape.has_text_frame if hasattr(shape, "has_text_frame") else False,
        "placeholder_type": None,
    }

    # Placeholder detection
    if shape.is_placeholder:
        info["placeholder_type"] = str(shape.placeholder_format.idx)
        if hasattr(shape, "placeholder_format") and shape.placeholder_format.type is not None:
            info["placeholder_type"] = str(shape.placeholder_format.type)

    # Text extraction
    if shape.has_text_frame:
        tf = shape.text_frame
        info["text"] = tf.text
        for para in tf.paragraphs:
            for run in para.runs:
                if run.font.name:
                    info["font_name"] = run.font.name
                if run.font.size:
                    info["font_size"] = run.font.size / 12700  # EMU → pt
                if run.font.bold is not None:
                    info["bold"] = run.font.bold
                if run.font.italic is not None:
                    info["italic"] = run.font.italic
                try:
                    if run.font.color and run.font.color.rgb:
                        info["color"] = str(run.font.color.rgb)
                except AttributeError:
                    # _NoneColor (inherited color) has no .rgb
                    pass
                break  # Take first run's style as representative
            break

    # Fill color (for shapes)
    if hasattr(shape, "fill"):
        try:
            fill = shape.fill
            if fill.type is not None:
                if hasattr(fill, "fore_color") and fill.fore_color.rgb:
                    info["fill_color"] = str(fill.fore_color.rgb)
        except Exception:
            pass

    # Image extraction
    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
        try:
            info["image_blob"] = shape.image.blob
        except Exception:
            pass

    # Table detection
    if shape.has_table:
        info["shape_type"] = "table"
        table = shape.table
        info["table_data"] = {
            "rows": table.rows.__len__(),
            "cols": len(table.columns),
            "content": [
                [cell.text for cell in row.cells]
                for row in table.rows
            ],
        }

    # Chart detection
    if shape.has_chart:
        info["shape_type"] = "chart"
        info["chart_type"] = str(shape.chart.chart_type) if hasattr(shape.chart, "chart_type") else "unknown"

    return info


def duplicate_presentation(source: Presentation) -> Presentation:
    """Deep-copy a Presentation, preserving all content and styles.

    Uses python-pptx's internal copy mechanism via saving to a BytesIO
    buffer and re-reading.

    Args:
        source: The source Presentation to copy.

    Returns:
        A new Presentation with identical content.
    """
    import io
    buffer = io.BytesIO()
    source.save(buffer)
    buffer.seek(0)
    return Presentation(buffer)
