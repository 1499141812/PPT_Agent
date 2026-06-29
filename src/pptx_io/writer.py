"""
PPTX writer — creates and saves presentations, applies edit operations.
"""

from __future__ import annotations

import copy
from pathlib import Path
from typing import Any, Optional

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor


def create_blank_presentation() -> Presentation:
    """Create a new blank Presentation.

    Returns:
        An empty ``pptx.Presentation``.
    """
    return Presentation()


def write_pptx(presentation: Presentation, output_path: str | Path) -> None:
    """Save a Presentation to disk.

    Args:
        presentation: The Presentation to save.
        output_path: Destination file path (should end with .pptx).
    """
    path = Path(output_path)
    path.parent.mkdir(parents=True, exist_ok=True)
    presentation.save(str(path))


def duplicate_slide(prs: Presentation, source_slide: Any) -> Any:
    """Clone an entire slide like PowerPoint's 'Duplicate Slide'.

    Copies the slide's spTree (shapes), background, and all image/media
    relationships so pictures render correctly on the new slide.
    """
    import copy
    from lxml import etree
    from pptx.oxml.ns import qn

    # 1. Create new slide with same layout
    new_slide = prs.slides.add_slide(source_slide.slide_layout)

    # 2. Copy all image/media relationships from source to new slide
    for rel in source_slide.part.rels.values():
        try:
            new_slide.part.relate_to(rel.target_part, rel.reltype)
        except Exception:
            pass

    # 3. Remove default shapes from new slide, copy source shapes
    spTree = new_slide.shapes._spTree
    for child in list(spTree):
        tag = child.tag.split('}')[-1] if '}' in child.tag else child.tag
        if tag in ('sp', 'pic', 'graphicFrame', 'grpSp', 'cxnSp'):
            spTree.remove(child)

    for child in source_slide.shapes._spTree:
        tag = child.tag.split('}')[-1] if '}' in child.tag else child.tag
        if tag in ('sp', 'pic', 'graphicFrame', 'grpSp', 'cxnSp'):
            spTree.append(copy.deepcopy(child))

    # 4. Copy per-slide background
    cSld_src = source_slide._element.find(qn('p:cSld'))
    cSld_dst = new_slide._element.find(qn('p:cSld'))
    if cSld_src is not None and cSld_dst is not None:
        bg = cSld_src.find(qn('p:bg'))
        if bg is not None:
            existing = cSld_dst.find(qn('p:bg'))
            if existing is not None:
                cSld_dst.remove(existing)
            cSld_dst.insert(0, copy.deepcopy(bg))

    return new_slide


def add_slide(
    presentation: Presentation,
    layout_index: int = 0,
) -> Any:
    """Add a new slide to the presentation.

    Args:
        presentation: Target Presentation.
        layout_index: Index of the slide layout to use (default: first / blank).

    Returns:
        The newly created Slide object.
    """
    layouts = presentation.slide_layouts
    if layout_index >= len(layouts):
        layout_index = 0
    slide_layout = layouts[layout_index]
    return presentation.slides.add_slide(slide_layout)


def clear_slide_text(slide: Any) -> dict[int, dict]:
    """Clear all text, returning saved font properties for restoration.

    Captures font SIZE and BOLD BEFORE tf.clear() (which destroys them).
    Returns dict mapping shape_id → {font_size_pt, bold}.
    """
    saved: dict[int, dict] = {}
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        tf = shape.text_frame
        # ── Capture BEFORE clearing ────────────────────────────────
        if len(tf.paragraphs) > 0:
            p0 = tf.paragraphs[0]
            r0 = p0.runs[0] if len(p0.runs) > 0 else None
            sz = r0.font.size if r0 else p0.font.size
            bold = r0.font.bold if r0 else p0.font.bold
            saved[shape.shape_id] = {
                "font_size": sz / 12700 if sz else None,  # EMU → pt
                "bold": bold,
            }
        tf.clear()
        if len(tf.paragraphs) == 0:
            tf.paragraphs.add()
    return saved


def delete_slide(presentation: Presentation, slide_index: int) -> None:
    """Remove a slide from the presentation by index.

    Uses the standard python-pptx XML manipulation approach.
    Reference: https://github.com/scanny/python-pptx/issues/67

    Args:
        presentation: Target Presentation.
        slide_index: 0-based index of the slide to delete.

    Raises:
        IndexError: If index is out of range.
    """
    if slide_index < 0 or slide_index >= len(presentation.slides):
        raise IndexError(f"Slide index {slide_index} out of range.")

    slides = presentation.slides._sldIdLst
    ns = "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}"
    rId = slides[slide_index].get(ns + "id")
    if rId is None:
        # Fallback: try the "id" attribute directly
        rId = slides[slide_index].get("id")

    # Remove the relationship (this also removes the slide part)
    if rId:
        try:
            presentation.part.drop_rel(rId)
        except Exception:
            pass  # Relationship might already be removed

    # Remove the slide entry from the slide list
    slides.remove(slides[slide_index])


# ── Shape manipulation helpers ──────────────────────────────────────────────

def _hex_to_rgb(hex_color: str) -> RGBColor:
    """Convert a hex color string (with or without #) to an RGBColor.

    Falls back to black (#000000) for any invalid input — the LLM sometimes
    produces malformed color strings like 'le' (truncated 'left' or 'blue').
    """
    import re
    hex_color = str(hex_color).strip().lstrip("#")
    # Validate: must be exactly 6 hex digits
    if not re.fullmatch(r"[0-9a-fA-F]{6}", hex_color):
        return RGBColor(0, 0, 0)  # safe fallback
    return RGBColor(
        int(hex_color[0:2], 16),
        int(hex_color[2:4], 16),
        int(hex_color[4:6], 16),
    )


def add_text_box(
    slide: Any,
    text: str,
    left: int | float,
    top: int | float,
    width: int | float,
    height: int | float,
    *,
    font_name: str = "Arial",
    font_size: int = 18,
    bold: bool = False,
    color: str = "#000000",
    alignment: str = "left",
) -> Any:
    """Add a text box to a slide.

    Args:
        slide: The target Slide.
        text: Text content.
        left, top, width, height: Position and size in EMU (or inches as float).
        font_name: Font family name.
        font_size: Font size in points.
        bold: Whether text is bold.
        color: Hex color string.
        alignment: "left", "center", or "right".

    Returns:
        The created Shape object.
    """
    # Convert inches → EMU if values are small (heuristic)
    if isinstance(left, float) and left < 100:
        left = int(left * 914400)
    if isinstance(top, float) and top < 100:
        top = int(top * 914400)
    if isinstance(width, float) and width < 100:
        width = int(width * 914400)
    if isinstance(height, float) and height < 100:
        height = int(height * 914400)

    txBox = slide.shapes.add_textbox(int(left), int(top), int(width), int(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.name = font_name
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = _hex_to_rgb(color)

    align_map = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT}
    p.alignment = align_map.get(alignment, PP_ALIGN.LEFT)

    return txBox


def add_image(
    slide: Any,
    image_path: str | Path,
    left: int | float,
    top: int | float,
    width: int | float | None = None,
    height: int | float | None = None,
) -> Any:
    """Add an image to a slide.

    Args:
        slide: The target Slide.
        image_path: Path to the image file (PNG, JPG, etc.).
        left, top: Position (inches as float, or EMU as int).
        width, height: Size in inches (float) or EMU (int). If None, natural size.

    Returns:
        The created Picture shape.
    """
    if isinstance(left, float) and left < 100:
        left = Inches(left)
    if isinstance(top, float) and top < 100:
        top = Inches(top)
    if width is not None and isinstance(width, float) and width < 100:
        width = Inches(width)
    if height is not None and isinstance(height, float) and height < 100:
        height = Inches(height)

    kwargs: dict[str, Any] = {}
    if width is not None:
        kwargs["width"] = width
    if height is not None:
        kwargs["height"] = height

    return slide.shapes.add_picture(
        str(Path(image_path).resolve()),
        int(left) if isinstance(left, int) else left,
        int(top) if isinstance(top, int) else top,
        **kwargs,
    )


def add_table(
    slide: Any,
    rows: int,
    cols: int,
    left: int | float,
    top: int | float,
    width: int | float,
    height: int | float,
    *,
    data: Optional[list[list[str]]] = None,
) -> Any:
    """Add a table to a slide.

    Args:
        slide: The target Slide.
        rows, cols: Table dimensions.
        left, top, width, height: Position and size.
        data: Optional 2D list of cell contents.

    Returns:
        The created Table shape.
    """
    if isinstance(left, float) and left < 100:
        left = int(left * 914400)
    if isinstance(top, float) and top < 100:
        top = int(top * 914400)
    if isinstance(width, float) and width < 100:
        width = int(width * 914400)
    if isinstance(height, float) and height < 100:
        height = int(height * 914400)

    table_shape = slide.shapes.add_table(
        rows, cols, int(left), int(top), int(width), int(height)
    )
    table = table_shape.table

    if data:
        for r in range(min(rows, len(data))):
            for c in range(min(cols, len(data[r]))):
                table.cell(r, c).text = str(data[r][c])

    return table_shape


def modify_text(shape: Any, new_text: str) -> None:
    """Change the text of a shape, PRESERVING font formatting.

    Instead of tf.clear() (which destroys font sizes/bold/color),
    replaces text only in the first paragraph's first run.
    """
    if not shape.has_text_frame:
        return
    tf = shape.text_frame

    # Clear text from all paragraphs but keep their formatting
    for para in tf.paragraphs:
        for run in para.runs:
            run.text = ""

    # Put new text in the first paragraph
    p0 = tf.paragraphs[0]
    if len(p0.runs) > 0:
        p0.runs[0].text = new_text
    else:
        p0.text = new_text


def delete_shape(slide: Any, shape: Any) -> None:
    """Remove a shape from a slide.

    Args:
        slide: The slide containing the shape.
        shape: The shape to remove.
    """
    sp = shape._element
    sp.getparent().remove(sp)

