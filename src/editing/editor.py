"""
Edit engine — executes edit operations against a real Presentation object.

This is the "executor" layer: it takes typed EditOperation objects and
applies them to python-pptx, translating inches to EMU, setting fonts,
inserting images, building tables, etc.
"""

from __future__ import annotations

import logging
from typing import Any, Optional

from pptx import Presentation

from src.editing.operations import (
    EditOperation,
    EditOpType,
    AddTextPayload,
    ModifyTextPayload,
    ModifyStylePayload,
    AddImagePayload,
    AddTablePayload,
    AddChartPayload,
    AddSlidePayload,
    DeleteShapePayload,
)
from src.pptx_io.writer import (
    add_text_box,
    add_image,
    add_table,
    add_slide,
    delete_slide,
    modify_text,
    delete_shape,
)

logger = logging.getLogger(__name__)


class EditEngine:
    """Applies edit operations to a python-pptx Presentation.

    Usage::

        engine = EditEngine(presentation)
        engine.execute(edit_op)
        # or
        engine.execute_batch([op1, op2, op3])
    """

    def __init__(self, presentation: Presentation) -> None:
        """Initialize with a target presentation.

        Args:
            presentation: The Presentation to edit (modified in-place).
        """
        self._prs = presentation
        self._executed_ops: list[EditOperation] = []

    # ── Public API ──────────────────────────────────────────────────────

    def execute(self, op: EditOperation) -> bool:
        """Execute a single edit operation.

        Args:
            op: The edit operation to apply.

        Returns:
            True if successful.

        Raises:
            ValueError: If the operation type is unknown.
            IndexError: If a slide index is out of range.
        """
        logger.debug("Executing: %s on slide %d", op.op_type.value, op.slide_idx)
        handler = self._DISPATCH.get(op.op_type)
        if handler is None:
            raise ValueError(f"Unknown operation type: {op.op_type}")

        handler(self, op)
        self._executed_ops.append(op)
        return True

    def execute_batch(self, ops: list[EditOperation]) -> int:
        """Execute a batch of edit operations.

        Args:
            ops: List of edit operations.

        Returns:
            Number of successful operations.

        Raises:
            No exception — individual failures are logged but don't stop the batch.
        """
        success_count = 0
        for op in ops:
            try:
                self.execute(op)
                success_count += 1
            except Exception as e:
                logger.warning("Edit failed (skipping): %s — %s", op.op_type.value, e)
        return success_count

    @property
    def history(self) -> list[EditOperation]:
        return list(self._executed_ops)

    @property
    def presentation(self) -> Presentation:
        return self._prs

    def _get_slide(self, idx: int) -> Any:
        while idx >= len(self._prs.slides):
            from src.pptx_io.writer import add_slide
            add_slide(self._prs)
        return self._prs.slides[idx]

    def _find_shape(self, slide: Any, shape_id=None, shape_name=None) -> Any:
        if shape_id is not None:
            for shape in slide.shapes:
                if shape.shape_id == shape_id:
                    return shape
        if shape_name is not None:
            for shape in slide.shapes:
                if shape.name == shape_name:
                    return shape
        return None

    def _handle_add_text(self, op): self._h_add_text(op)
    def _handle_modify_text(self, op): self._h_modify_text(op)
    def _handle_modify_style(self, op): self._h_modify_style(op)
    def _handle_add_image(self, op): self._h_add_image(op)
    def _handle_add_table(self, op): self._h_add_table(op)
    def _handle_add_chart(self, op): self._h_add_chart(op)
    def _handle_add_slide(self, op): self._h_add_slide(op)
    def _handle_delete_slide(self, op): self._h_delete_slide(op)
    def _handle_delete_shape(self, op): self._h_delete_shape(op)
    def _handle_reorder_slide(self, op): self._h_reorder_slide(op)

    _DISPATCH = {
        EditOpType.ADD_TEXT: _handle_add_text,
        EditOpType.MODIFY_TEXT: _handle_modify_text,
        EditOpType.MODIFY_STYLE: _handle_modify_style,
        EditOpType.ADD_IMAGE: _handle_add_image,
        EditOpType.ADD_TABLE: _handle_add_table,
        EditOpType.ADD_CHART: _handle_add_chart,
        EditOpType.ADD_SLIDE: _handle_add_slide,
        EditOpType.DELETE_SLIDE: _handle_delete_slide,
        EditOpType.DELETE_SHAPE: _handle_delete_shape,
        EditOpType.REORDER_SLIDE: _handle_reorder_slide,
    }

    def _h_add_text(self, op):
        payload = _ensure_payload(op, AddTextPayload)
        slide = self._get_slide(op.slide_idx)
        from src.pptx_io.writer import add_text_box
        add_text_box(slide, text=payload.text,
                     left=payload.position.left, top=payload.position.top,
                     width=payload.position.width, height=payload.position.height,
                     font_name=payload.style.font_name, font_size=int(payload.style.font_size),
                     bold=payload.style.bold, color=payload.style.color,
                     alignment=payload.style.alignment)

    def _h_modify_text(self, op):
        payload = _ensure_payload(op, ModifyTextPayload)
        slide = self._get_slide(op.slide_idx)
        shape = self._find_shape(slide, payload.shape_id, payload.shape_name)
        if shape is None:
            raise ValueError(f"modify_text: shape id={payload.shape_id} not found")
        from src.pptx_io.writer import modify_text
        modify_text(shape, payload.new_text)

    def _h_modify_style(self, op):
        payload = _ensure_payload(op, ModifyStylePayload)
        slide = self._get_slide(op.slide_idx)
        shape = self._find_shape(slide, payload.shape_id, payload.shape_name)
        if shape is None:
            raise ValueError(f"modify_style: shape id={payload.shape_id} not found")
        if not shape.has_text_frame:
            raise ValueError(f"modify_style: shape has no text frame")
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                if "font_name" in payload.changes:
                    run.font.name = payload.changes["font_name"]
                if "font_size" in payload.changes:
                    from pptx.util import Pt
                    run.font.size = Pt(payload.changes["font_size"])
                if "bold" in payload.changes:
                    run.font.bold = payload.changes["bold"]
                if "color" in payload.changes:
                    from src.pptx_io.writer import _hex_to_rgb
                    run.font.color.rgb = _hex_to_rgb(payload.changes["color"])
                break

    def _h_add_image(self, op):
        payload = _ensure_payload(op, AddImagePayload)
        from pathlib import Path
        if not Path(payload.image_path).exists():
            raise FileNotFoundError(f"Image not found: {payload.image_path}")
        slide = self._get_slide(op.slide_idx)
        from src.pptx_io.writer import add_image
        add_image(slide, image_path=payload.image_path,
                  left=payload.position.left, top=payload.position.top,
                  width=payload.position.width or None, height=payload.position.height or None)

    def _h_add_table(self, op):
        payload = _ensure_payload(op, AddTablePayload)
        slide = self._get_slide(op.slide_idx)
        from src.pptx_io.writer import add_table
        add_table(slide, rows=len(payload.rows)+1, cols=len(payload.headers),
                  left=payload.position.left, top=payload.position.top,
                  width=payload.position.width, height=payload.position.height,
                  data=[payload.headers] + payload.rows)

    def _h_add_chart(self, op):
        payload = _ensure_payload(op, AddChartPayload)
        from src.generation.chart_generator import ChartGenerator
        gen = ChartGenerator()
        img_path = gen.generate_chart(chart_type=payload.chart_type, data=payload.data, title=payload.title)
        slide = self._get_slide(op.slide_idx)
        from src.pptx_io.writer import add_image
        add_image(slide, image_path=img_path,
                  left=payload.position.left, top=payload.position.top,
                  width=payload.position.width, height=payload.position.height)

    def _h_add_slide(self, op):
        from src.pptx_io.writer import add_slide
        payload = _ensure_payload(op, AddSlidePayload)
        add_slide(self._prs, layout_index=payload.layout_index)

    def _h_delete_slide(self, op):
        from src.pptx_io.writer import delete_slide
        delete_slide(self._prs, op.slide_idx)

    def _h_delete_shape(self, op):
        payload = _ensure_payload(op, DeleteShapePayload)
        slide = self._get_slide(op.slide_idx)
        shape = self._find_shape(slide, payload.shape_id, payload.shape_name)
        if shape is None:
            raise ValueError(f"delete_shape: shape id={payload.shape_id} not found")
        from src.pptx_io.writer import delete_shape
        delete_shape(slide, shape)

    def _h_reorder_slide(self, op):
        from_idx = op.payload.get("from_index", 0) if isinstance(op.payload, dict) else op.payload.from_index
        to_idx = op.payload.get("to_index", 0) if isinstance(op.payload, dict) else op.payload.to_index
        if from_idx == to_idx:
            return
        sldIdLst = self._prs.slides._sldIdLst
        elements = list(sldIdLst)
        if 0 <= from_idx < len(elements) and 0 <= to_idx < len(elements):
            elem = elements.pop(from_idx)
            elements.insert(to_idx, elem)
            for e in elements:
                sldIdLst.append(e)


def _ensure_payload(op: EditOperation, model_cls: type) -> Any:
    if isinstance(op.payload, model_cls):
        return op.payload
    if isinstance(op.payload, dict):
        try:
            return model_cls(**op.payload)
        except Exception as e:
            raise TypeError(
                f"Payload mismatch for {op.op_type.value}: "
                f"{model_cls.__name__} expects specific fields, "
                f"but got keys {list(op.payload.keys())}. Detail: {e}"
            ) from e
    raise TypeError(
        f"Expected {model_cls.__name__} but payload is {type(op.payload).__name__}."
    )

def _ensure_payload(op: EditOperation, model_cls: type) -> Any:
    """Ensure the payload matches the expected type for this op_type.

    LLMs sometimes confuse payload structures (e.g. op_type="add_image"
    with modify_style-shaped payload). We detect mismatches and raise
    clear errors so execute_batch can skip them gracefully.
    """
    if isinstance(op.payload, model_cls):
        return op.payload
    if isinstance(op.payload, dict):
        try:
            return model_cls(**op.payload)
        except Exception as e:
            raise TypeError(
                f"Payload mismatch for {op.op_type.value}: "
                f"{model_cls.__name__} expects specific fields, "
                f"but got keys {list(op.payload.keys())}. "
                f"Detail: {e}"
            ) from e
    raise TypeError(
        f"Expected {model_cls.__name__} but payload is {type(op.payload).__name__}. "
        f"LLM likely confused op_type with wrong payload shape."
    )

