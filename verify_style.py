"""
Standalone verification: extract style from reference PPT → apply to new PPT.

Usage: python verify_style.py reference.pptx [output.pptx]
"""

import sys
from pathlib import Path

def main():
    ref_path = Path(sys.argv[1]) if len(sys.argv) > 1 else None
    out_path = Path(sys.argv[2]) if len(sys.argv) > 2 else Path("verify_output.pptx")

    if not ref_path or not ref_path.exists():
        print("Usage: python verify_style.py reference.pptx [output.pptx]")
        sys.exit(1)

    from pptx import Presentation
    from pptx.oxml.ns import qn
    from lxml import etree
    import copy

    ref = Presentation(str(ref_path))
    print(f"Reference: {len(ref.slides)} slides, {len(ref.slide_layouts)} layouts")

    # Pick slide 0 as representative
    ref_slide = ref.slides[0]

    # ── 0. Dump ALL shapes for diagnosis ────────────────────────────────
    print(f"0. All shapes on reference slide 0:")
    for i, s in enumerate(ref_slide.shapes):
        stype_raw = s.shape_type
        stype = str(stype_raw) if stype_raw else 'None'
        classname = type(s).__name__
        has_img = hasattr(s, 'image') and s.image is not None
        text = s.text_frame.text[:50] if s.has_text_frame else ''
        print(f"   [{i}] type={stype}, class={classname}, "
              f"has_image={has_img}, text='{text}'")

    print()

    # ── 1. Extract background ──────────────────────────────────────────
    cSld = ref_slide._element.find(qn('p:cSld'))
    bg = cSld.find(qn('p:bg')) if cSld is not None else None
    bg_xml = etree.tostring(bg, encoding='unicode') if bg is not None else None
    print(f"1. Background extracted: {bg is not None}")
    if bg is not None:
        bgPr = bg.find(qn('p:bgPr'))
        solid = bgPr.find(qn('a:solidFill')) if bgPr is not None else None
        grad = bgPr.find(qn('a:gradFill')) if bgPr is not None else None
        print(f"   Type: {'solid' if solid is not None else 'gradient' if grad is not None else 'other'}")

    # ── 2. Extract fonts ───────────────────────────────────────────────
    from src.graph.workflow import _extract_slide_fonts
    fonts = _extract_slide_fonts(ref_slide)
    print(f"2. Fonts extracted: main={fonts['main_font']}, "
          f"title_size={fonts['title_size']}, body_size={fonts['body_size']}, "
          f"color=#{fonts['main_color']}")
    print(f"   All fonts: {fonts['all_fonts']}")
    print(f"   All colors: {fonts['all_colors']}")

    # ── 3. Extract shapes ──────────────────────────────────────────────
    from src.graph.workflow import _extract_slide_shapes
    shapes_xml = _extract_slide_shapes(ref_slide, filter_content=True)
    print(f"3. Shapes extracted: {len(shapes_xml)} decorative shapes")

    # Check shape types and dump fills
    pic_count = sum(1 for s in shapes_xml if s.get("type") == "picture")
    xml_count = sum(1 for s in shapes_xml if s.get("type") == "xml")
    print(f"   Pictures: {pic_count}, XML shapes: {xml_count}")
    for i, s in enumerate(shapes_xml):
        if s.get("type") == "xml":
            xml = s["xml"]
            has_solid = 'solidFill' in xml
            has_grad = 'gradFill' in xml
            has_blip = 'blipFill' in xml
            has_scheme = 'schemeClr' in xml
            has_srgb = 'srgbClr' in xml
            print(f"   XML[{i}]: solidFill={has_solid}, gradFill={has_grad}, "
                  f"blipFill={has_blip}, schemeClr={has_scheme}, srgbClr={has_srgb}")

    # ── 4. Create new presentation and apply ────────────────────────────
    from src.pptx_io.reader import duplicate_presentation
    from src.pptx_io.writer import delete_slide, add_text_box

    new = duplicate_presentation(ref)
    while len(new.slides) > 0:
        delete_slide(new, 0)
    print(f"4. New PPT created: {len(new.slides)} slides, {len(new.slide_layouts)} layouts")

    # Add a slide with layout 0
    new_slide = new.slides.add_slide(new.slide_layouts[0])
    print(f"   Added slide with layout 0")

    # ── 5. Apply background ────────────────────────────────────────────
    if bg_xml:
        from src.graph.workflow import _apply_slide_background
        _apply_slide_background(new_slide, bg_xml)
        # Verify
        cSld2 = new_slide._element.find(qn('p:cSld'))
        bg2 = cSld2.find(qn('p:bg')) if cSld2 is not None else None
        print(f"5. BG applied: {bg2 is not None}")
    else:
        print(f"5. BG skipped (none in reference)")

    # ── 6. Inject decorative shapes ─────────────────────────────────────
    from src.graph.workflow import _inject_shapes
    _inject_shapes(new_slide, shapes_xml)
    print(f"6. Shapes injected: {len(new_slide.shapes)} total shapes on slide")

    # ── 7. Add text with enforced font ──────────────────────────────────
    add_text_box(new_slide, "Test Title 测试标题", left=1, top=0.5, width=8, height=1.2,
                 font_name=fonts['main_font'], font_size=fonts['title_size'],
                 bold=True, color=fonts['main_color'])
    add_text_box(new_slide, "Body text 正文内容", left=1, top=2, width=8, height=3,
                 font_name=fonts['main_font'], font_size=fonts['body_size'],
                 color=fonts['main_color'])

    # Apply enforcement
    from src.graph.workflow import _enforce_font_style
    _enforce_font_style(new_slide, fonts)
    print(f"7. Font style enforced: font={fonts['main_font']}, color=#{fonts['main_color']}")

    # ── 8. Save and verify ──────────────────────────────────────────────
    new.save(str(out_path))
    print(f"8. Saved to: {out_path}")

    # Reopen and verify
    verify = Presentation(str(out_path))
    vs = verify.slides[0]

    # Check background
    vcSld = vs._element.find(qn('p:cSld'))
    vbg = vcSld.find(qn('p:bg')) if vcSld is not None else None
    print(f"\n   VERIFY: BG present = {vbg is not None}")

    # Check shapes
    print(f"   VERIFY: {len(vs.shapes)} shapes on slide")
    for s in vs.shapes:
        stype = str(s.shape_type) if s.shape_type else '?'
        text = s.text_frame.text[:40] if s.has_text_frame else '(no text)'
        print(f"     [{stype}] '{text}'")

    print(f"\nDone. Open {out_path} in PowerPoint to check.")

if __name__ == "__main__":
    main()
